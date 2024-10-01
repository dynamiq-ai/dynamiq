import copy
from io import BytesIO
from pathlib import Path
from typing import Any, Literal

from pptx import Presentation

from dynamiq.components.converters.base import BaseConverter
from dynamiq.components.converters.utils import get_filename_for_bytesio
from dynamiq.types import Document, DocumentCreationMode


class PPTXConverter(BaseConverter):
    """
    A component for converting files to Documents using the pptx library.

          Initializes the object with the configuration for converting documents using
          the python-pptx.

          Args:
          document_creation_mode (Literal["one-doc-per-file", "one-doc-per-page", "one-doc-per-element"], optional):
              Determines how to create Documents from the elements of presentation. Options are:
              - `"one-doc-per-file"`: Creates one Document per file.
                  All elements are concatenated into one text field.
              - `"one-doc-per-page"`: Creates one Document per page.
                  All elements on a page are concatenated into one text field.
              Defaults to `"one-doc-per-file"`.
          extraction_mode(ExtractionMode): Type of text extraction format.

        Usage example:
            ```python
            from dynamiq.components.converters.pptx import PPTXConverter

            converter = PPTXConverter()
            documents = converter.run(paths=["a/file/path.pptx", "a/directory/path"])["documents"]
            ```
    """

    document_creation_mode: Literal[DocumentCreationMode.ONE_DOC_PER_FILE, DocumentCreationMode.ONE_DOC_PER_PAGE] = (
        DocumentCreationMode.ONE_DOC_PER_FILE
    )

    def _process_file(self, file: Path | BytesIO, metadata: dict[str, Any]) -> list[Any]:
        """
        Process a single presentation and create documents.

        Args:
            file (Union[Path, BytesIO]): The file to process.
            metadata (Dict[str, Any]): Metadata to attach to the documents.

        Returns:
            List[Any]: A list of created documents.

        Raises:
            ValueError: If the file object doesn't have a name and its extension can't be guessed.
            TypeError: If the file argument is neither a Path nor a BytesIO object.
        """
        if isinstance(file, Path):
            with open(file, "rb") as upload_file:
                file_content = BytesIO(upload_file.read())
                file_path = upload_file.name
        elif isinstance(file, BytesIO):
            file_path = get_filename_for_bytesio(file)
            file_content = file
        else:
            raise TypeError("Expected a Path object or a BytesIO object.")
        elements = Presentation(file_content)
        return self._create_documents(
            filepath=file_path,
            elements=elements,
            document_creation_mode=self.document_creation_mode,
            metadata=metadata,
        )

    def _create_documents(
        self,
        filepath: str,
        elements: Any,
        document_creation_mode: DocumentCreationMode,
        metadata: dict[str, Any],
        **kwargs
    ) -> list[Document]:
        """
        Create Documents from the elements of the presentation.
        """
        docs = []
        if document_creation_mode == DocumentCreationMode.ONE_DOC_PER_FILE:
            text_all_slides = "\n".join(
                "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text)
                for slide in elements.slides
            )
            metadata = copy.deepcopy(metadata)
            metadata["file_path"] = filepath
            docs = [Document(content=text_all_slides, metadata=metadata)]

        elif document_creation_mode == DocumentCreationMode.ONE_DOC_PER_PAGE:
            texts_per_page = [
                "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text)
                for slide in elements.slides
            ]
            metadata = copy.deepcopy(metadata)
            metadata["file_path"] = filepath

            docs = [Document(content=text, metadata=metadata) for text in texts_per_page]

        return docs
