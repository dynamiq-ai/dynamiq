"""
HTML to PPTX Converter

Converts HTML slides to PowerPoint presentations with formatting preservation.
"""

import re

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


class HTMLToPPTXConverter:
    """Converts HTML slides to PPTX format."""

    def __init__(self):
        self.prs = None
        self.slide_width = Inches(10)  # 16:9 ratio default
        self.slide_height = Inches(5.625)

    def convert(self, html_path, output_path, slide_width=960, slide_height=540, preserve_formatting=True):
        """
        Convert HTML file to PPTX presentation.

        Args:
            html_path: Path to HTML file
            output_path: Path for output PPTX file
            slide_width: Slide width in pixels (default 960 for 16:9)
            slide_height: Slide height in pixels (default 540 for 16:9)
            preserve_formatting: Preserve CSS formatting

        Returns:
            Path to created PPTX file
        """
        # Read HTML
        with open(html_path, encoding="utf-8") as f:
            html_content = f.read()

        # Parse HTML
        soup = BeautifulSoup(html_content, 'html.parser')

        # Create presentation
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width / 96)  # Convert pixels to inches (96 DPI)
        self.prs.slide_height = Inches(slide_height / 96)

        # Find all slide divs
        slides = soup.find_all('div', class_='slide')

        for slide_div in slides:
            self._create_slide_from_html(slide_div, preserve_formatting)

        # Save presentation
        self.prs.save(output_path)
        return output_path

    def _create_slide_from_html(self, slide_div, preserve_formatting):
        """Create a PowerPoint slide from HTML div."""
        # Add blank slide
        blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Extract background color
        bg_color = self._extract_css_property(slide_div, 'background', '#FFFFFF')
        if bg_color and bg_color != 'transparent':
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self._parse_color(bg_color)

        # Process child elements
        self._process_elements(slide, slide_div, preserve_formatting)

    def _process_elements(self, slide, parent_element, preserve_formatting, y_offset=0):
        """Process HTML elements and add to slide."""
        current_y = y_offset

        for element in parent_element.children:
            if element.name is None:  # Text node
                continue

            # Extract styles
            styles = self._extract_styles(element)

            if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Add title text box
                text = element.get_text().strip()
                if text:
                    left = self._parse_size(styles.get('margin-left', '60px'))
                    top = current_y if current_y > 0 else self._parse_size(styles.get('margin-top', '60px'))
                    width = self._parse_size(styles.get('width', '800px'))
                    height = self._parse_size(styles.get('height', '80px'))

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text

                    # Apply formatting
                    if preserve_formatting:
                        p = text_frame.paragraphs[0]
                        font = p.font
                        font.size = self._parse_font_size(styles.get('font-size', '48pt'))
                        font.bold = 'bold' in styles.get('font-weight', '')
                        color = styles.get('color', '#000000')
                        if color:
                            font.color.rgb = self._parse_color(color)

                    current_y = top + height + Inches(0.2)

            elif element.name in ['p', 'div']:
                # Add paragraph
                text = element.get_text().strip()
                if text and not element.find_all(['h1', 'h2', 'h3', 'ul', 'ol']):  # Skip containers
                    left = self._parse_size(styles.get('margin-left', '60px'))
                    top = current_y if current_y > 0 else self._parse_size(styles.get('margin-top', '120px'))
                    width = self._parse_size(styles.get('width', '800px'))
                    height = Inches(1)

                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True

                    if preserve_formatting:
                        p = text_frame.paragraphs[0]
                        font = p.font
                        font.size = self._parse_font_size(styles.get('font-size', '24pt'))
                        color = styles.get('color', '#000000')
                        if color:
                            font.color.rgb = self._parse_color(color)

                    current_y = top + height + Inches(0.1)

            elif element.name in ['ul', 'ol']:
                # Add bullet points
                items = element.find_all('li')
                left = self._parse_size(styles.get('margin-left', '100px'))
                top = current_y if current_y > 0 else Inches(2)
                width = self._parse_size(styles.get('width', '700px'))
                height = Inches(len(items) * 0.5)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame

                for idx, item in enumerate(items):
                    item_text = item.get_text().strip()
                    p = text_frame.add_paragraph() if idx > 0 else text_frame.paragraphs[0]
                    p.text = item_text
                    p.level = 0

                    if preserve_formatting:
                        font = p.font
                        font.size = self._parse_font_size(styles.get('font-size', '24pt'))
                        color = styles.get('color', '#000000')
                        if color:
                            font.color.rgb = self._parse_color(color)

                current_y = top + height + Inches(0.2)

            # Recursively process nested elements
            if element.name == 'div' and element.find_all(['h1', 'h2', 'h3', 'p', 'ul']):
                self._process_elements(slide, element, preserve_formatting, current_y)

    def _extract_styles(self, element):
        """Extract CSS styles from element."""
        styles = {}
        style_attr = element.get('style', '')

        if style_attr:
            for style in style_attr.split(';'):
                if ':' in style:
                    key, value = style.split(':', 1)
                    styles[key.strip()] = value.strip()

        # Also check class-based styles (simplified)
        return styles

    def _extract_css_property(self, element, property_name, default=''):
        """Extract CSS property value."""
        style = element.get('style', '')
        match = re.search(rf'{property_name}:\s*([^;]+)', style)
        return match.group(1).strip() if match else default

    def _parse_color(self, color_str):
        """Parse color string to RGBColor."""
        # Handle hex colors
        if color_str.startswith('#'):
            hex_color = color_str.lstrip('#')
            if len(hex_color) == 3:
                hex_color = ''.join([c*2 for c in hex_color])
            r, g, b = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
            return RGBColor(r, g, b)

        # Handle rgb() format
        rgb_match = re.match(r'rgb\((\d+),\s*(\d+),\s*(\d+)\)', color_str)
        if rgb_match:
            r, g, b = map(int, rgb_match.groups())
            return RGBColor(r, g, b)

        # Default to black
        return RGBColor(0, 0, 0)

    def _parse_font_size(self, size_str):
        """Parse font size string to Pt."""
        # Extract number from size string
        match = re.match(r'(\d+)', size_str)
        if match:
            size = int(match.group(1))
            return Pt(size)
        return Pt(24)  # Default

    def _parse_size(self, size_str):
        """Parse size string to Inches."""
        # Extract number
        match = re.match(r'(\d+)', size_str)
        if match:
            pixels = int(match.group(1))
            return Inches(pixels / 96)  # Convert pixels to inches
        return Inches(1)


def generate_thumbnails(pptx_path, output_dir):
    """
    Generate thumbnail images for each slide.

    Args:
        pptx_path: Path to PPTX file
        output_dir: Directory to save thumbnails

    Returns:
        List of thumbnail paths
    """
    import os

    from PIL import Image, ImageDraw
    from pptx import Presentation

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)

    thumbnails = []
    for idx, slide in enumerate(prs.slides):
        # Create blank image
        img = Image.new('RGB', (960, 540), color='white')
        draw = ImageDraw.Draw(img)

        # Draw placeholder (actual rendering would need more complex logic)
        draw.text((50, 50), f"Slide {idx + 1}", fill='black')

        # Save thumbnail
        thumb_path = os.path.join(output_dir, f"slide_{idx + 1}.png")
        img.save(thumb_path)
        thumbnails.append(thumb_path)

    return thumbnails


# Example usage
if __name__ == "__main__":
    converter = HTMLToPPTXConverter()
    converter.convert(
        html_path="presentation.html",
        output_path="presentation.pptx",
        slide_width=960,
        slide_height=540
    )
    print("Conversion complete!")
