from io import BytesIO
from pathlib import Path

import pytest
from pptx import Presentation

from dynamiq import Workflow
from dynamiq.flows import Flow
from dynamiq.nodes.converters.pptx import PPTXFileConverter, PPTXFileConverterInputSchema
from dynamiq.nodes.node import NodeDependency
from dynamiq.nodes.utils import Output
from dynamiq.runnables import RunnableResult, RunnableStatus
from dynamiq.types import Document


def write_presentation_to_path(path: Path, presentation: Presentation) -> str:
    presentation.save(path)
    return str(path)


def write_presentation_to_bytesio(presentation: Presentation, filename: str = "file.pptx") -> BytesIO:
    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    buffer.name = filename
    return buffer


def create_pptx_with_content(content: str) -> tuple[Presentation, str]:
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = content
    return prs, content


def create_empty_pptx() -> Presentation:
    return Presentation()


@pytest.fixture
def pptx_converter():
    return PPTXFileConverter(
        id="pptx_converter",
        name="Test PPTX Converter",
    )


@pytest.fixture
def output_node(pptx_converter):
    return Output(id="output_node", depends=[NodeDependency(pptx_converter)])


@pytest.fixture
def workflow_with_pptx_converter_and_output(pptx_converter, output_node):
    return Workflow(
        id="test_workflow",
        flow=Flow(
            nodes=[pptx_converter, output_node],
        ),
        version="1",
    )


@pytest.fixture
def valid_pptx_content():
    return "Hello, World!"


@pytest.fixture
def valid_pptx_presentation(valid_pptx_content):
    return create_pptx_with_content(valid_pptx_content)


@pytest.fixture
def valid_pptx_file_path(tmp_path, valid_pptx_presentation):
    prs, content = valid_pptx_presentation
    file_path = tmp_path / "valid.pptx"
    return write_presentation_to_path(file_path, prs), content


@pytest.fixture
def valid_pptx_bytesio(valid_pptx_presentation):
    prs, content = valid_pptx_presentation
    file = write_presentation_to_bytesio(prs, "valid.pptx")
    return file, content


@pytest.fixture
def empty_pptx_presentation_obj():
    return create_empty_pptx()


@pytest.fixture
def empty_pptx_presentation_file_path(tmp_path, empty_pptx_presentation_obj):
    file_path = tmp_path / "empty_presentation.pptx"
    return write_presentation_to_path(file_path, empty_pptx_presentation_obj)


@pytest.fixture
def empty_pptx_presentation_bytesio(empty_pptx_presentation_obj):
    return write_presentation_to_bytesio(empty_pptx_presentation_obj, "empty_presentation.pptx")


@pytest.fixture
def empty_pptx_file(tmp_path):
    empty_file_path = tmp_path / "empty.pptx"
    empty_file_path.touch()
    return str(empty_file_path)


@pytest.fixture
def invalid_pptx_content():
    return b"This is not a valid PPTX file content"


@pytest.fixture
def invalid_pptx_file_path(tmp_path, invalid_pptx_content):
    invalid_file = tmp_path / "invalid.pptx"
    invalid_file.write_bytes(invalid_pptx_content)
    return str(invalid_file)


@pytest.fixture
def invalid_pptx_bytesio(invalid_pptx_content):
    file = BytesIO(invalid_pptx_content)
    file.name = "invalid.pptx"
    return file


@pytest.fixture
def unsupported_content():
    return b"This is not a PPTX file, just plain text"


@pytest.fixture
def unsupported_file_path(tmp_path, unsupported_content):
    text_file = tmp_path / "text.txt"
    text_file.write_bytes(unsupported_content)
    return str(text_file)


@pytest.fixture
def unsupported_bytesio(unsupported_content):
    wrong_file = BytesIO(unsupported_content)
    wrong_file.name = "text.txt"
    return wrong_file


@pytest.fixture
def non_existent_pptx_file(tmp_path):
    return str(tmp_path / "non_existent_file.pptx")


@pytest.mark.parametrize(
    "input_type,input_fixture",
    [
        ("file_paths", "valid_pptx_file_path"),
        ("files", "valid_pptx_bytesio"),
    ],
)
def test_workflow_with_pptx_converter(request, input_type, input_fixture):
    file_or_path, content = request.getfixturevalue(input_fixture)
    pptx_converter = PPTXFileConverter()
    wf_pptx = Workflow(flow=Flow(nodes=[pptx_converter]))
    input_data = {input_type: [file_or_path]}

    response = wf_pptx.run(input_data=input_data)
    document_id = response.output[next(iter(response.output))]["output"]["documents"][0]["id"]

    expected_file_path = file_or_path if input_type == "file_paths" else file_or_path.name

    pptx_converter_expected_result = RunnableResult(
        status=RunnableStatus.SUCCESS,
        input=dict(PPTXFileConverterInputSchema(**input_data)),
        output={"documents": [Document(id=document_id, content=content, metadata={"file_path": expected_file_path})]},
    ).to_dict(skip_format_types={BytesIO, bytes})

    expected_output = {pptx_converter.id: pptx_converter_expected_result}
    assert response == RunnableResult(status=RunnableStatus.SUCCESS, input=input_data, output=expected_output)


@pytest.mark.parametrize(
    "input_type,input_fixture",
    [
        ("file_paths", "invalid_pptx_file_path"),
        ("files", "invalid_pptx_bytesio"),
    ],
)
def test_workflow_with_pptx_converter_parsing_error(
    request, workflow_with_pptx_converter_and_output, pptx_converter, output_node, input_type, input_fixture
):
    invalid_input = request.getfixturevalue(input_fixture)
    input_data = {input_type: [invalid_input]}

    response = workflow_with_pptx_converter_and_output.run(input_data=input_data)

    assert response.status == RunnableStatus.FAILURE
    assert response.output[pptx_converter.id]["status"] == RunnableStatus.FAILURE.value
    assert "error" in response.output[pptx_converter.id]
    assert response.output[output_node.id]["status"] == RunnableStatus.SKIP.value


def test_workflow_with_pptx_converter_file_not_found(
    workflow_with_pptx_converter_and_output, pptx_converter, output_node, non_existent_pptx_file
):
    input_data = {"file_paths": [non_existent_pptx_file]}

    response = workflow_with_pptx_converter_and_output.run(input_data=input_data)

    assert response.status == RunnableStatus.FAILURE
    assert response.output[pptx_converter.id]["status"] == RunnableStatus.FAILURE.value
    assert "No files found in the provided paths" in response.output[pptx_converter.id]["error"]["message"]
    assert response.output[output_node.id]["status"] == RunnableStatus.SKIP.value


def test_workflow_with_pptx_converter_empty_file(
    workflow_with_pptx_converter_and_output, pptx_converter, output_node, empty_pptx_file
):
    input_data = {"file_paths": [empty_pptx_file]}

    response = workflow_with_pptx_converter_and_output.run(input_data=input_data)

    assert response.status == RunnableStatus.FAILURE
    assert response.output[pptx_converter.id]["status"] == RunnableStatus.FAILURE.value
    assert "error" in response.output[pptx_converter.id]
    assert response.output[output_node.id]["status"] == RunnableStatus.SKIP.value


@pytest.mark.parametrize(
    "input_type,input_fixture",
    [
        ("file_paths", "empty_pptx_presentation_file_path"),
        ("files", "empty_pptx_presentation_bytesio"),
    ],
)
def test_workflow_with_pptx_converter_empty_presentation(
    request, workflow_with_pptx_converter_and_output, pptx_converter, output_node, input_type, input_fixture
):
    empty_pptx = request.getfixturevalue(input_fixture)
    input_data = {input_type: [empty_pptx]}
    response = workflow_with_pptx_converter_and_output.run(input_data=input_data)

    assert response.status == RunnableStatus.SUCCESS
    assert response.output[pptx_converter.id]["status"] == RunnableStatus.SUCCESS.value
    assert "documents" in response.output[pptx_converter.id]["output"]
    documents = response.output[pptx_converter.id]["output"]["documents"]
    assert len(documents) == 1
    assert documents[0]["content"] == ""

    expected_path = empty_pptx if input_type == "file_paths" else empty_pptx.name
    assert documents[0]["metadata"]["file_path"] == expected_path

    assert response.output[output_node.id]["status"] == RunnableStatus.SUCCESS.value


@pytest.mark.parametrize(
    "input_type,input_fixture",
    [
        ("file_paths", "unsupported_file_path"),
        ("files", "unsupported_bytesio"),
    ],
)
def test_workflow_with_pptx_converter_unsupported_file(
    request, workflow_with_pptx_converter_and_output, pptx_converter, output_node, input_type, input_fixture
):
    unsupported = request.getfixturevalue(input_fixture)
    input_data = {input_type: [unsupported]}

    response = workflow_with_pptx_converter_and_output.run(input_data=input_data)

    assert response.status == RunnableStatus.FAILURE
    assert response.output[pptx_converter.id]["status"] == RunnableStatus.FAILURE.value
    assert "error" in response.output[pptx_converter.id]
    assert response.output[output_node.id]["status"] == RunnableStatus.SKIP.value
